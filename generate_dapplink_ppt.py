import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 主题色
BLUE = RGBColor(0, 112, 192)
GREEN = RGBColor(0, 176, 80)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(10.33)
    height = Inches(2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "DappLink 简介"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BLUE
    tf.word_wrap = True
    p.alignment = PP_ALIGN.CENTER
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = "构建 Web3 中间件，赋能区块链未来"
    p2.font.size = Pt(28)
    p2.font.color.rgb = GREEN
    p2.alignment = PP_ALIGN.CENTER
    # 英文副标题
    p3 = tf.add_paragraph()
    p3.text = "Building Web3 Middleware to Empower the Future of Blockchain"
    p3.font.size = Pt(18)
    p3.font.color.rgb = BLACK
    p3.alignment = PP_ALIGN.CENTER

def add_catalog_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(10.33)
    height = Inches(5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = "目录"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    for i, item in enumerate(["中文简介", "产品与服务", "合作伙伴", "英文简介", "社交账号"], 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.level = 1

def add_content_slide(title, content, color=BLUE, subcontent=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    # 内容
    left = Inches(1.2)
    top = Inches(2)
    width = Inches(11)
    height = Inches(4.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf2 = content_box.text_frame
    for line in content:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
    if subcontent:
        for line in subcontent:
            p = tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = BLUE
            p.level = 1

def add_social_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "社交账号 / Social Media"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN
    # 内容
    left = Inches(1.2)
    top = Inches(2)
    width = Inches(11)
    height = Inches(4.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf2 = content_box.text_frame
    accounts = [
        ("官网 / Website", "https://www.dapplink.xyz/zh"),
        ("推特 / Twitter", "https://x.com/0xdapplink"),
        ("Github", "https://github.com/dapplink-labs"),
        ("Telegram", "https://t.me/+qqhy1i-_xnU1M2Jl"),
        ("Medium", "https://medium.com/@0xappslink"),
    ]
    for name, url in accounts:
        p = tf2.add_paragraph()
        p.text = f"{name}: {url}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

add_title_slide()
add_catalog_slide()
add_content_slide(
    "中文简介",
    [
        "DappLink 正在重塑 Web3 应用的开发方式，让构建区块链项目变得更简单、更高效，助力用户快速实现产品落地。",
        "我们提供一揽子的 Web3 中间件技术解决方案。"
    ],
    color=BLUE
)
add_content_slide(
    "产品与服务",
    [
        "• 一键部署钱包（含中心化钱包、去中心化钱包和托管系统）",
        "• 一键部署可组合 Dapp",
        "• 一键部署应用链（Layer2、Layer3 和 Cosmos 应用链）",
        "• 一键部署 Layer2 快速验证网络",
        "• 一站式 RWA & PayFi 技术支持"
    ],
    color=GREEN
)
add_content_slide(
    "合作伙伴",
    [
        "已与多个头部生态及项目达成深度合作，包括：",
        "Manta、CpChain、DataMining、Parapack、FishCake，",
        "以及多家位于迪拜的 Layer2 和 Cosmos 的应用链项目。",
        "持续推动 Web3 技术的落地与普及。"
    ],
    color=BLUE
)
add_content_slide(
    "英文简介",
    [
        "DappLink is redefining how Web3 applications are developed—making blockchain project deployment simpler and more efficient, helping users quickly bring their products to market.",
        "We offer a comprehensive suite of Web3 middleware solutions."
    ],
    color=GREEN
)
add_social_slide()

prs.save("DappLink_企业招商简介.pptx")
print("PPT已生成：DappLink_企业招商简介.pptx") 